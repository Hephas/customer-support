import os
import io
import json
import base64
import google.generativeai as genai
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2 import service_account
from http.server import BaseHTTPRequestHandler

# ── PDF
import fitz  # PyMuPDF

# ── Word
from docx import Document

# ── Excel
import openpyxl

# ── PowerPoint
from pptx import Presentation

# ══════════════════════════════════════════
# 群翌能源 System Prompt
# ══════════════════════════════════════════
SYSTEM_PROMPT = """你是群翌能源（Hephas Energy）的專業客服AI助理。

## 核心角色
- 代表群翌能源提供專業、親切的客戶服務
- 協助客戶解決產品諮詢、技術問題及售後服務需求
- 維護公司專業形象，提升客戶滿意度

## 公司基本資訊
- 公司：群翌能源股份有限公司（Hephas Energy Corporation）
- 專業：氫能源設備、燃料電池測試設備、關鍵系統零組件
- 地址：台灣新竹縣新竹科學園區園區二路60號1F
- 電話：+886-3-578-0221
- 官網：https://www.hephasenergy.com
- Email：info@hephasenergy.com

## 回應規範
- 必須全程使用繁體中文
- 語氣專業、有禮貌、親切
- 優先使用提供的文件資料回答
- 文件中找不到答案時，誠實告知並建議聯繫專人

## 禁止事項
- 不可編造技術數據或產品規格
- 不可承諾無法確認的事項

## 轉人工客服
需要時請說：「建議聯繫專人客服 📞 +886-3-578-0221 / 📧 info@hephasenergy.com」
"""

# ══════════════════════════════════════════
# 初始化設定
# ══════════════════════════════════════════
genai.configure(api_key=os.environ.get("GEMINI_API_KEY"))

DRIVE_FOLDER_ID = os.environ.get("GOOGLE_DRIVE_FOLDER_ID", "")
MAX_FILES       = 3      # 每次最多查幾個檔案
MAX_CHARS       = 4000   # 每個檔案最多取幾個字

# ══════════════════════════════════════════
# Google Drive 服務帳號初始化
# ══════════════════════════════════════════
def get_drive_service():
    """從環境變數取得 Service Account，建立 Drive 連線"""
    key_b64 = os.environ.get("GOOGLE_SERVICE_ACCOUNT_KEY_B64", "")
    if not key_b64:
        return None
    try:
        key_json = json.loads(base64.b64decode(key_b64).decode("utf-8"))
        creds = service_account.Credentials.from_service_account_info(
            key_json,
            scopes=["https://www.googleapis.com/auth/drive.readonly"]
        )
        return build("drive", "v3", credentials=creds)
    except Exception as e:
        print(f"[Drive Init Error] {e}")
        return None

# ══════════════════════════════════════════
# 搜尋相關文件
# ══════════════════════════════════════════
def search_relevant_files(service, query):
    """用關鍵字在指定資料夾搜尋最相關的文件"""
    try:
        results = service.files().list(
            q=f"'{DRIVE_FOLDER_ID}' in parents and trashed=false",
            fields="files(id, name, mimeType)",
            pageSize=50
        ).execute()

        files = results.get("files", [])
        keywords = [k for k in query.lower().split() if len(k) > 1]

        scored = []
        for f in files:
            name_lower = f["name"].lower()
            score = sum(1 for kw in keywords if kw in name_lower)
            scored.append((score, f))

        scored.sort(key=lambda x: -x[0])

        # 如果完全沒有關鍵字命中，還是回傳前幾個當背景知識
        top = [f for _, f in scored[:MAX_FILES]]
        return top

    except Exception as e:
        print(f"[Search Error] {e}")
        return []

# ══════════════════════════════════════════
# 各格式文字提取
# ══════════════════════════════════════════
def extract_text(service, file_info):
    """根據 MIME 類型提取文件文字"""
    mime    = file_info["mimeType"]
    file_id = file_info["id"]
    name    = file_info["name"]
    text    = ""

    try:
        # ── Google 原生格式（直接匯出）
        if mime == "application/vnd.google-apps.document":
            raw  = service.files().export(fileId=file_id, mimeType="text/plain").execute()
            text = raw.decode("utf-8")

        elif mime == "application/vnd.google-apps.spreadsheet":
            raw  = service.files().export(fileId=file_id, mimeType="text/csv").execute()
            text = raw.decode("utf-8")

        elif mime == "application/vnd.google-apps.presentation":
            raw  = service.files().export(fileId=file_id, mimeType="text/plain").execute()
            text = raw.decode("utf-8")

        else:
            # ── 下載二進制檔案
            buf = io.BytesIO()
            req = service.files().get_media(fileId=file_id)
            dl  = MediaIoBaseDownload(buf, req)
            done = False
            while not done:
                _, done = dl.next_chunk()
            buf.seek(0)

            # ── PDF
            if mime == "application/pdf":
                doc  = fitz.open(stream=buf.read(), filetype="pdf")
                text = "\n".join(page.get_text() for page in doc)

            # ── Word
            elif mime in [
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword"
            ]:
                doc  = Document(buf)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

            # ── Excel
            elif mime in [
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel"
            ]:
                wb   = openpyxl.load_workbook(buf, data_only=True)
                rows = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    rows.append(f"[工作表：{sheet_name}]")
                    for row in ws.iter_rows(values_only=True):
                        line = " | ".join(str(c) for c in row if c is not None)
                        if line.strip():
                            rows.append(line)
                text = "\n".join(rows)

            # ── PowerPoint
            elif mime in [
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint"
            ]:
                prs  = Presentation(buf)
                slides = []
                for i, slide in enumerate(prs.slides, 1):
                    slides.append(f"[第 {i} 頁]")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slides.append(shape.text)
                text = "\n".join(slides)

    except Exception as e:
        print(f"[Extract Error] {name}: {e}")
        text = f"（{name} 讀取失敗）"

    return f"📄 【{name}】\n{text[:MAX_CHARS]}"

# ══════════════════════════════════════════
# 主要 HTTP Handler
# ══════════════════════════════════════════
class handler(BaseHTTPRequestHandler):

    def do_OPTIONS(self):
        self.send_response(200)
        self._cors()
        self.end_headers()

    def do_POST(self):
        try:
            length  = int(self.headers.get("Content-Length", 0))
            body    = json.loads(self.rfile.read(length))
            user_msg = body.get("message", "").strip()
            history  = body.get("history", [])

            if not user_msg:
                self._json(400, {"error": "訊息不可為空"})
                return

            # ── 嘗試查詢 Google Drive
            drive_context = ""
            drive_service = get_drive_service()

            if drive_service and DRIVE_FOLDER_ID:
                files = search_relevant_files(drive_service, user_msg)
                if files:
                    parts = [extract_text(drive_service, f) for f in files]
                    drive_context = "\n\n".join(parts)

            # ── 組合最終 Prompt
            final_system = SYSTEM_PROMPT
            if drive_context:
                final_system += f"""

## 📚 參考文件（來自公司資料庫）
以下是從公司文件中找到的相關資料，請優先根據這些內容回答：

{drive_context}

---
回答時可說「根據我們的產品資料...」
"""

            # ── 對話歷史
            chat_history = []
            for item in history[:-1]:
                role = "user" if item.get("role") == "user" else "model"
                chat_history.append({"role": role, "parts": [item.get("content", "")]})

            # ── 呼叫 Gemini
            model = genai.GenerativeModel(
                model_name="gemini-1.5-flash",
                system_instruction=final_system
            )
            chat  = model.start_chat(history=chat_history)
            reply = chat.send_message(user_msg).text

            self._json(200, {"reply": reply})

        except json.JSONDecodeError:
            self._json(400, {"error": "無效的請求格式"})
        except Exception as e:
            print(f"[Server Error] {e}")
            self._json(500, {"error": "伺服器錯誤，請稍後再試"})

    def _json(self, code, data):
        self.send_response(code)
        self._cors()
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False).encode("utf-8"))

    def _set_cors_headers(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")

    def _cors(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
